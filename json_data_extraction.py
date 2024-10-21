import streamlit as st
from PyPDF2 import PdfReader
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
import google.generativeai as genai
from docx import Document
from pptx import Presentation
import pandas as pd
import os
import json
from dotenv import load_dotenv

load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=api_key)

if not api_key:
    st.error("API key not found. Please check your .env file or environment variables.")

class PolicyExtractor:
    def __init__(self):
        self.model = genai.GenerativeModel(model_name='gemini-1.5-pro')

    def get_text_from_document(self, uploaded_files):
        text = ""
        for file in uploaded_files:
            if file.name.lower().endswith((".pdf", ".docx", ".txt", ".pptx", ".csv", ".json")):
                try:
                    if file.name.lower().endswith(".pdf"):
                        pdf_reader = PdfReader(file)
                        for page in pdf_reader.pages:
                            text += page.extract_text() or ""
                    elif file.name.lower().endswith(".docx"):
                        doc_content = Document(file)
                        text += "\n".join([paragraph.text for paragraph in doc_content.paragraphs])
                    elif file.name.lower().endswith(".txt"):
                        text += file.read().decode('utf-8', errors='ignore')
                    elif file.name.lower().endswith(".pptx"):
                        presentation = Presentation(file)
                        for slide in presentation.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                    elif file.name.lower().endswith('.csv'):
                        df = pd.read_csv(file)
                        text += df.to_csv(index=False)
                    elif file.name.lower().endswith('.json'):
                        json_data = json.load(file)
                        text  += json.dumps(json_data, indent=4) 
                except Exception as e:
                    st.error(f"Error processing {file.name}: {e}")
        return text
    
    def get_text_chunks(self, text):
        """Split text into chunks."""
        if not text:
            print("No text provided.")
            return []
        try:
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
            chunks = text_splitter.split_text(text)
            return chunks
        except Exception as e:
            print(f"Error splitting text: {e}")
            return []

    def extract_policy_information(self, chunks):
        """Extract specific insurance policy information from the uploaded document content."""
        prompt = """
        Extract all relevant information from the uploaded property damage document and return it in the following structured JSON format. 

        Follow these rules:
        1. If any field is missing, leave it empty.
        2. Capture all nested data accurately.
        3. Ensure numerical fields are correctly parsed (e.g., year, phone number, policy limits).
        4. For 'categories', 'subcategories', and 'questions maintain the structure as shown.

        Use the following JSON template:

        {
            "inspections": [
                {
                    "id": 1,
                    "name": "",
                    "claim_information": {
                        "claim_number": "",
                        "claim_type_id": "",
                        "created_by_id": "",
                        "customer_id": "",
                        "id": "",
                        "line_of_business_id": "",
                        "status": "",
                        "synched": 0
                    },
                    "property_information": {
                        "claim_id": "",
                        "construction_type": "",
                        "created_at": "",
                        "created_by_id": "",
                        "id": "",
                        "location_id": "",
                        "number_of_stories": 0,
                        "owner_name": "",
                        "property_type": "",
                        "roof_type": "",
                        "synched": 0,
                        "year_of_built": 0
                    },
                    "customer_information": {
                        "created_at": "",
                        "created_by_id": "",
                        "email": "",
                        "first_name": "",
                        "id": "",
                        "last_name": "",
                        "location_address": "",
                        "location_city": "",
                        "location_id": "",
                        "location_pincode": "",
                        "location_state": "",
                        "phone_number": "",
                        "updated_at": ""
                    },
                    "insurance_policy_information": {
                        "claim_id": "",
                        "coverage_type": "",
                        "created_by_id": "",
                        "deductible": "",
                        "id": "",
                        "insurance_carrier": "",
                        "policy_holder_name": "",
                        "policy_limits": "",
                        "policy_number": "",
                        "synched": 0
                    },
                    "damage_information": {
                        "cause_of_loss": "",
                        "claim_id": "",
                        "created_by_id": "",
                        "date_of_loss": "",
                        "description": "",
                        "id": "",
                        "synched": 0
                    },
                    "claim_type": {
                        "id": 1,
                        "name": ""
                    },
                    "categories": [
                        {
                        "id": 1
                        "name": "",
                        "is_interior": false,
                        "priority": 0,
                        "subcategories": [
                            {
                            "id": 1,
                            "title": "",
                            "description": "",
                            "helptext": "",
                            "priority": 1,
                            "questions": [
                                {
                                "id": 1,
                                "title": "",
                                "description": "",
                                "helptext": "",
                                "priority": 1,
                                "required": true
                                "answer_type": "",
                                "photos": true,
                                "videos": true,
                                "photos_360": true,
                                "photos_response_collection": [
                                    {}
                                ],
                                "docs": false,
                                "video_response_collection": [],
                                "360_photo_response_collection": [],
                                "notes": true,
                                "applicable": true,
                                "is_additional_questions": false
                                "response": ""
                                }
                            ]
                            }
                        ]
                        }
                    ]
                    }
                }
            ]
        }
        You do not need to adhere to any predefined format. Simply use the document structure and content to organize the output appropriately.
        Please extract relevant data from the document into this structure.
        """
        extracted_info = {}
        try:
            combined_input = f"{prompt}\n\n" + '\n'.join(chunks)
            input_data = {
                "parts": [
                    {"text": combined_input}
                ]
            }
            response = self.model.generate_content(input_data) 

            if response.candidates:
                generated_content = response.candidates[0].content.parts[0].text.strip() 
                generated_content = generated_content.replace("```json", "").replace("```", "").strip() 
                try:
                    extracted_info = json.loads(generated_content)   
                except json.JSONDecodeError as e:
                    extracted_info = {"error": "Invalid JSON response", "content": generated_content}
            else:
                extracted_info = {"error": "No candidates found in response."}

        except Exception as e:
            print("Error during extraction:", e)
            extracted_info = {"error": str(e)}

        return extracted_info


    def main(self):
        """Main function to run the policy extraction application."""
        st.set_page_config("Insurance Policy JSON Generator")
        st.header("Generate JSON from Insurance Policy Documents:")
        uploaded_files = st.file_uploader("Upload Policy Documents", type=["pdf", "docx", "txt", "pptx", "csv", "json"], accept_multiple_files=True)

        if uploaded_files:
            st.write("Generating JSON...")
            raw_text = self.get_text_from_document(uploaded_files)
            chunks = self.get_text_chunks(raw_text)
            policy_info_dict = self.extract_policy_information(chunks)
            policy_info_json = json.dumps(policy_info_dict, indent=4)
            st.write("Generated JSON:")
            st.code(policy_info_json, language='json')
            # st.json(policy_info_dict)
            st.download_button(
                label="Download JSON",
                data=policy_info_json,
                file_name='policy_info.json',
                mime='application/json'
            )

if __name__ == "__main__":
    app = PolicyExtractor()
    app.main()
